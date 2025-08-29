from flask import Flask, request, render_template
import pandas as pd
import os
import platform
import subprocess
import sys
from datetime import datetime
from pathlib import Path
from collections import defaultdict
import webbrowser
import threading
import logging
from pptx import Presentation


# Detectar sistema operativo
ON_WINDOWS = platform.system() == "Windows"

# Configurar logging
logging.basicConfig(
    level=logging.DEBUG,
    filename="app.log",
    filemode="w",
    format="%(asctime)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(__name__)

# Configuración para PyInstaller
def resource_path(relative_path):
    """Obtener ruta absoluta de recursos, funciona tanto en desarrollo como en .exe"""
    try:
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

app = Flask(__name__,
           template_folder=resource_path('templates'),
           static_folder=resource_path('static'))

def open_browser():
    try:
        webbrowser.open_new("http://127.0.0.1:5000/")
    except Exception as e:
        logger.error(f"Error abriendo navegador: {e}")

def get_downloads_folder():
    try:
        home = Path.home()
        downloads = home / "Downloads"
        if not downloads.exists():
            downloads = home / "Descargas"
            if not downloads.exists():
                downloads = home
        return downloads
    except Exception as e:
        logger.error(f"Error obteniendo carpeta de descargas: {e}")
        return Path.cwd()

def get_plantilla_path_pptx():
    possible = [
        resource_path("plantilla.pptx"),
        "plantilla.pptx",
        Path.cwd() / "plantilla.pptx",
        get_downloads_folder() / "certificados" / "plantilla.pptx"
    ]
    for p in possible:
        if os.path.exists(p):
            logger.info(f"Plantilla PPTX encontrada en: {p}")
            return str(p)
    raise FileNotFoundError("No se encontró plantilla.pptx en ninguna ubicación")

# funciones para manejo de PPTX y PDF

def safe_filename(s: str) -> str:
    keep = (" ", ".", "_", "-")
    return "".join(c for c in s if c.isalnum() or c in keep).rstrip()

def build_placeholder_map(context: dict) -> dict:
    mapping = {}
    for k, v in context.items():
        vstr = "" if v is None else str(v)
        mapping[f"{{{{{k.upper()}}}}}"] = vstr
        mapping[f"{{{{{k.lower()}}}}}"] = vstr
        mapping[k.upper()] = vstr
        mapping[k.lower()] = vstr
    return mapping

def replace_placeholders_in_presentation(prs: Presentation, mapping: dict):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            try:
                tf = shape.text_frame
            except Exception:
                continue

            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    orig = run.text or ""
                    new = orig
                    for ph, val in mapping.items():
                        if ph in new:
                            new = new.replace(ph, val)
                    if new != orig:
                        run.text = new

def render_pptx_template(template_path: str, context: dict, out_pptx_path: str):
    logger.info(f"Renderizando PPTX: {template_path} -> {out_pptx_path}")
    prs = Presentation(template_path)
    mapping = build_placeholder_map(context)
    replace_placeholders_in_presentation(prs, mapping)
    prs.save(out_pptx_path)


def convert_pptx_to_pdf_powerpoint_fixed(pptx_path: str, pdf_path: str) -> bool:
    """
    Versión mejorada de conversión con PowerPoint COM
    """
    if not ON_WINDOWS:
        logger.error("Conversión con PowerPoint solo disponible en Windows.")
        return False

    try:
        import win32com.client
        import pythoncom
        
        logger.info(f"Iniciando conversión con PowerPoint: {pptx_path} -> {pdf_path}")

        # Inicializar COM
        pythoncom.CoInitialize()
        
        # Rutas absolutas
        pptx_path_abs = os.path.abspath(pptx_path)
        pdf_path_abs = os.path.abspath(pdf_path)
        
        if not os.path.exists(pptx_path_abs):
            logger.error(f"Archivo PPTX no encontrado: {pptx_path_abs}")
            return False
        
        # Crear aplicación PowerPoint
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        
        # Minimizar ventana en lugar de ocultarla completamente
        ppt_app.Visible = 1
        ppt_app.WindowState = 2  # Minimized
        
        logger.info(f"Abriendo presentación: {pptx_path_abs}")
        presentation = ppt_app.Presentations.Open(
            pptx_path_abs,
            ReadOnly=1,
            Untitled=1,
            WithWindow=0
        )
        
        logger.info(f"Exportando a PDF: {pdf_path_abs}")
        
        # Usar ExportAsFixedFormat con parámetros optimizados
        presentation.ExportAsFixedFormat(
            pdf_path_abs,
            2,  # ppFixedFormatTypePDF
            Intent=1,  # ppFixedFormatIntentPrint - Para impresión de calidad
            FrameSlides=0,  # No enmarcar slides
            HandoutOrder=1,
            OutputType=2,  # ppPrintOutputSlides
            PrintHiddenSlides=0,
            PrintRange=None,
            RangeType=1,  # ppPrintAll
            SlideShowName="",
            IncludeDocProps=1,
            KeepIRMSettings=1,
            DocStructureTags=1,
            BitmapMissingFonts=1,
            UseDocumentICCProfile=0
        )
        
        # Cerrar y limpiar
        presentation.Close()
        ppt_app.Quit()
        
        # Liberar recursos COM
        del presentation
        del ppt_app
        pythoncom.CoUninitialize()
        
        # Verificar resultado
        if os.path.exists(pdf_path_abs) and os.path.getsize(pdf_path_abs) > 1000:  # Al menos 1KB
            logger.info(f"PDF generado exitosamente: {pdf_path_abs} ({os.path.getsize(pdf_path_abs)} bytes)")
            return True
        else:
            logger.error("PowerPoint no generó PDF válido")
            return False

    except Exception as e:
        logger.error(f"Error en conversión PowerPoint: {str(e)}")
        try:
            if 'presentation' in locals():
                presentation.Close()
            if 'ppt_app' in locals():
                ppt_app.Quit()
            pythoncom.CoUninitialize()
        except:
            pass
        return False


def convert_pptx_to_pdf_advanced_python(pptx_path: str, pdf_path: str) -> bool:
    """
    Conversión avanzada usando python-pptx + reportlab
    Intenta mantener el formato visual lo mejor posible
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter, A4, landscape
        from reportlab.lib.units import inch, mm
        from reportlab.lib.colors import HexColor, black, navy, darkblue
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        import io
        
        logger.info(f"Conversión avanzada Python: {pptx_path} -> {pdf_path}")
        
        # Cargar presentación
        prs = Presentation(pptx_path)
        
        # Usar orientación landscape para certificados (más común)
        page_width, page_height = landscape(A4)
        
        # Crear canvas
        c = canvas.Canvas(pdf_path, pagesize=landscape(A4))
        
        # Configurar fuentes
        try:
            # Intentar usar fuentes del sistema
            c.setFont("Helvetica-Bold", 24)
        except:
            c.setFont("Helvetica", 24)
        
        # Procesar cada slide
        for slide_num, slide in enumerate(prs.slides):
            if slide_num > 0:
                c.showPage()  # Nueva página
            
            logger.info(f"Procesando slide {slide_num + 1}")
            
            # Obtener dimensiones del slide
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Factor de escala para ajustar al PDF
            scale_x = page_width / (slide_width / 914400)  # EMU to points conversion
            scale_y = page_height / (slide_height / 914400)
            scale = min(scale_x, scale_y) * 0.9  # 90% para margen
            
            # Centrar contenido
            offset_x = (page_width - (slide_width / 914400) * scale) / 2
            offset_y = (page_height - (slide_height / 914400) * scale) / 2
            
            # Analizar y colocar texto por posición
            text_elements = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Obtener posición y tamaño del shape
                    left = shape.left / 914400 * scale + offset_x  # EMU to points
                    top = page_height - (shape.top / 914400 * scale + offset_y)
                    width = shape.width / 914400 * scale
                    height = shape.height / 914400 * scale
                    
                    text_elements.append({
                        'text': shape.text.strip(),
                        'left': left,
                        'top': top,
                        'width': width,
                        'height': height
                    })
            
            # Ordenar elementos por posición vertical (top to bottom)
            text_elements.sort(key=lambda x: -x['top'])
            
            # Dibujar elementos
            for i, element in enumerate(text_elements):
                text = element['text']
                x = element['left']
                y = element['top']
                
                # Determinar estilo basado en posición y contenido
                if i == 0 or 'INSTITUTO' in text.upper() or 'CAMPUSLANDS' in text.upper():
                    # Título principal
                    c.setFont("Helvetica-Bold", 16)
                    c.setFillColor(darkblue)
                elif 'HACE CONSTAR' in text.upper():
                    # Subtítulo
                    c.setFont("Helvetica-Bold", 14)
                    c.setFillColor(black)
                elif any(placeholder in text for placeholder in ['{{NOMBRE}}', '{{CEDULA}}', '{{HORAS}}']):
                    # Texto con datos variables
                    c.setFont("Helvetica-Bold", 12)
                    c.setFillColor(navy)
                elif 'RECTOR' in text.upper() or 'COORDINADOR' in text.upper():
                    # Firmas
                    c.setFont("Helvetica-Bold", 10)
                    c.setFillColor(black)
                else:
                    # Texto normal
                    c.setFont("Helvetica", 11)
                    c.setFillColor(black)
                
                # Manejar texto multilínea
                lines = text.split('\n')
                line_height = 15
                
                for j, line in enumerate(lines):
                    if line.strip():
                        # Centrar texto si parece ser un título
                        if (i < 3 or 'INSTITUTO' in line.upper() or 
                            'CAMPUSLANDS' in line.upper() or 'HACE CONSTAR' in line.upper()):
                            text_width = c.stringWidth(line, "Helvetica-Bold", 16)
                            x_centered = (page_width - text_width) / 2
                            c.drawString(x_centered, y - (j * line_height), line)
                        else:
                            c.drawString(x, y - (j * line_height), line)
            
            # Agregar bordes decorativos si es un certificado
            if slide_num == 0:  # Solo en el primer slide
                # Borde exterior
                c.setStrokeColor(darkblue)
                c.setLineWidth(3)
                margin = 20
                c.rect(margin, margin, page_width - 2*margin, page_height - 2*margin)
                
                # Borde interior
                c.setStrokeColor(navy)
                c.setLineWidth(1)
                inner_margin = 30
                c.rect(inner_margin, inner_margin, 
                      page_width - 2*inner_margin, page_height - 2*inner_margin)
        
        # Finalizar PDF
        c.save()
        
        if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 1000:
            logger.info(f"PDF avanzado generado: {pdf_path} ({os.path.getsize(pdf_path)} bytes)")
            return True
        else:
            logger.error("No se generó PDF válido con método avanzado")
            return False
            
    except ImportError as e:
        logger.error(f"Librerías no disponibles: {e}")
        return False
    except Exception as e:
        logger.error(f"Error en conversión avanzada Python: {str(e)}")
        return False


def convert_pptx_to_pdf_libreoffice(pptx_path: str, output_dir: str) -> bool:

    """
    Conversión usando LibreOffice (multiplataforma)
    """
    try:
        pptx_path_abs = os.path.abspath(pptx_path)
        output_dir_abs = os.path.abspath(output_dir)
        
        # Posibles ubicaciones de LibreOffice en Windows
        possible_soffice = [
            "soffice",
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        
        soffice_path = None
        for path in possible_soffice:
            if os.path.exists(path) or path == "soffice":
                try:
                    # Probar si funciona
                    result = subprocess.run([path, "--version"], 
                                          capture_output=True, text=True, timeout=10)
                    if result.returncode == 0:
                        soffice_path = path
                        break
                except:
                    continue
        
        if not soffice_path:
            logger.error("LibreOffice no encontrado")
            return False
        
        cmd = [
            soffice_path,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", output_dir_abs,
            pptx_path_abs
        ]
        
        logger.info(f"Ejecutando LibreOffice: {' '.join(cmd)}")
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        
        if result.returncode == 0:
            base_name = os.path.splitext(os.path.basename(pptx_path))[0]
            expected_pdf = os.path.join(output_dir_abs, f"{base_name}.pdf")
            
            if os.path.exists(expected_pdf) and os.path.getsize(expected_pdf) > 0:
                logger.info(f"PDF generado con LibreOffice: {expected_pdf}")
                return True
                
        logger.error(f"LibreOffice falló: {result.stderr}")
        return False
            
    except Exception as e:
        logger.error(f"Error con LibreOffice: {str(e)}")
        return False


def convert_pptx_to_pdf_with_preview(pptx_path: str, pdf_path: str) -> bool:
    """
    Método usando pillow para convertir slides a imágenes y luego a PDF
    Mantiene exactamente el formato visual
    """
    try:
        # Este método requiere que PowerPoint genere imágenes primero
        import win32com.client
        from PIL import Image
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import landscape, A4
        from reportlab.lib.utils import ImageReader
        import tempfile
        import shutil
        
        logger.info(f"Conversión con preview: {pptx_path} -> {pdf_path}")
        
        # Crear directorio temporal
        temp_dir = tempfile.mkdtemp()
        
        try:
            # Abrir PowerPoint y exportar como imágenes
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            ppt_app.Visible = 1
            ppt_app.WindowState = 2  # Minimized
            
            presentation = ppt_app.Presentations.Open(os.path.abspath(pptx_path))
            
            # Exportar cada slide como imagen
            for i in range(1, presentation.Slides.Count + 1):
                img_path = os.path.join(temp_dir, f"slide_{i}.png")
                presentation.Slides(i).Export(img_path, "PNG", 1920, 1080)  # Alta resolución
            
            presentation.Close()
            ppt_app.Quit()
            
            # Crear PDF desde las imágenes
            c = canvas.Canvas(pdf_path, pagesize=landscape(A4))
            page_width, page_height = landscape(A4)
            
            for i in range(1, presentation.Slides.Count + 1):
                if i > 1:
                    c.showPage()
                
                img_path = os.path.join(temp_dir, f"slide_{i}.png")
                if os.path.exists(img_path):
                    # Abrir imagen y ajustar al tamaño de página
                    img = Image.open(img_path)
                    img_width, img_height = img.size
                    
                    # Calcular escala manteniendo proporción
                    scale_x = (page_width - 40) / img_width  # 40 puntos de margen
                    scale_y = (page_height - 40) / img_height
                    scale = min(scale_x, scale_y)
                    
                    new_width = img_width * scale
                    new_height = img_height * scale
                    
                    # Centrar imagen
                    x = (page_width - new_width) / 2
                    y = (page_height - new_height) / 2
                    
                    # Dibujar imagen
                    c.drawImage(ImageReader(img), x, y, new_width, new_height)
            
            c.save()
            
            # Limpiar archivos temporales
            shutil.rmtree(temp_dir)
            
            if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 1000:
                logger.info(f"PDF con preview generado: {pdf_path}")
                return True
                
        except Exception as e:
            logger.error(f"Error en conversión con preview: {e}")
            shutil.rmtree(temp_dir, ignore_errors=True)
            return False
            
    except ImportError:
        logger.error("Librerías para preview no disponibles (pillow)")
        return False
    except Exception as e:
        logger.error(f"Error general en preview: {e}")
        return False



def convert_pptx_to_pdf_python_libs(pptx_path: str, pdf_path: str) -> bool:
    """
    Conversión usando librerías de Python (python-pptx + reportlab)
    Este método es más básico pero no requiere software adicional
    """
    try:
        from pptx import Presentation
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_CENTER, TA_LEFT
        
        logger.info(f"Convirtiendo a PDF con librerías Python: {pptx_path} -> {pdf_path}")
        
        # Cargar presentación
        prs = Presentation(pptx_path)
        
        # Crear PDF
        doc = SimpleDocTemplate(pdf_path, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # Estilo personalizado para certificados
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=24,
            spaceAfter=30,
            alignment=TA_CENTER,
            textColor='navy'
        )
        
        normal_style = ParagraphStyle(
            'CustomNormal',
            parent=styles['Normal'],
            fontSize=14,
            spaceAfter=12,
            alignment=TA_CENTER
        )
        
        # Procesar cada slide
        for slide_num, slide in enumerate(prs.slides):
            if slide_num > 0:
                story.append(Spacer(1, 0.5*inch))
                
            # Extraer texto de cada shape
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            # Agregar contenido al PDF
            if slide_text:
                # El primer texto como título
                if slide_text[0]:
                    story.append(Paragraph(slide_text[0], title_style))
                
                # Resto como párrafos normales
                for text in slide_text[1:]:
                    if text:
                        story.append(Paragraph(text, normal_style))
                        story.append(Spacer(1, 0.2*inch))
        
        # Generar PDF
        doc.build(story)
        
        if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
            logger.info(f"PDF generado con librerías Python: {pdf_path}")
            return True
        else:
            logger.error("No se pudo generar PDF con librerías Python")
            return False
            
    except ImportError as e:
        logger.error(f"Librerías no disponibles para conversión Python: {e}")
        logger.error("Instala con: pip install reportlab")
        return False
    except Exception as e:
        logger.error(f"Error en conversión con librerías Python: {str(e)}")
        return False


def convert_pptx_to_pdf_robust(pptx_path: str, pdf_path: str) -> bool:
    """
    Función principal que intenta todos los métodos disponibles
    """
    logger.info(f"Iniciando conversión robusta: {pptx_path} -> {pdf_path}")
    
    # Método 1: PowerPoint COM (Windows con Office)
    if ON_WINDOWS:
        logger.info("Intentando conversión con PowerPoint COM...")
        if convert_pptx_to_pdf_powerpoint(pptx_path, pdf_path):
            return True
        logger.warning("PowerPoint COM falló, probando siguiente método...")
    
    # Método 2: LibreOffice
    logger.info("Intentando conversión con LibreOffice...")
    output_dir = os.path.dirname(pdf_path)
    if convert_pptx_to_pdf_libreoffice(pptx_path, output_dir):
        # LibreOffice genera el PDF con nombre base del PPTX
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        generated_pdf = os.path.join(output_dir, f"{base_name}.pdf")
        
        if generated_pdf != pdf_path and os.path.exists(generated_pdf):
            try:
                if os.path.exists(pdf_path):
                    os.remove(pdf_path)
                os.rename(generated_pdf, pdf_path)
                logger.info(f"PDF renombrado: {generated_pdf} -> {pdf_path}")
            except Exception as e:
                logger.error(f"Error renombrando PDF: {e}")
                return False
        return True
    
    logger.warning("LibreOffice falló, probando último método...")
    
    # Método 3: Librerías Python (básico pero funcional)
    logger.info("Intentando conversión con librerías Python...")
    if convert_pptx_to_pdf_python_libs(pptx_path, pdf_path):
        return True
    
    logger.error("Todos los métodos de conversión fallaron")
    return False


def convert_pptx_to_pdf_ultimate(pptx_path: str, pdf_path: str) -> bool:
    """
    Función principal que intenta todos los métodos, priorizando calidad visual
    """
    logger.info(f"=== CONVERSIÓN DEFINITIVA ===")
    logger.info(f"PPTX: {pptx_path}")
    logger.info(f"PDF: {pdf_path}")
    
    # Método 1: PowerPoint COM mejorado (mejor calidad)
    if ON_WINDOWS:
        logger.info("🔄 Método 1: PowerPoint COM mejorado")
        if convert_pptx_to_pdf_powerpoint_fixed(pptx_path, pdf_path):
            logger.info("✅ PowerPoint COM exitoso")
            return True
        logger.warning("❌ PowerPoint COM falló")
    
    # Método 2: PowerPoint con preview (calidad perfecta)
    if ON_WINDOWS:
        logger.info("🔄 Método 2: PowerPoint con preview")
        if convert_pptx_to_pdf_with_preview(pptx_path, pdf_path):
            logger.info("✅ PowerPoint preview exitoso")
            return True
        logger.warning("❌ PowerPoint preview falló")
    
    # Método 3: LibreOffice (buena calidad, multiplataforma)
    logger.info("🔄 Método 3: LibreOffice")
    output_dir = os.path.dirname(pdf_path)
    if convert_pptx_to_pdf_libreoffice(pptx_path, output_dir):
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        generated_pdf = os.path.join(output_dir, f"{base_name}.pdf")
        
        if generated_pdf != pdf_path and os.path.exists(generated_pdf):
            try:
                if os.path.exists(pdf_path):
                    os.remove(pdf_path)
                os.rename(generated_pdf, pdf_path)
                logger.info("✅ LibreOffice exitoso")
                return True
            except Exception as e:
                logger.error(f"Error renombrando: {e}")
        elif os.path.exists(generated_pdf):
            logger.info("✅ LibreOffice exitoso")
            return True
    logger.warning("❌ LibreOffice falló")
    
    # Método 4: Python avanzado (respaldo mejorado)
    logger.info("🔄 Método 4: Python avanzado")
    if convert_pptx_to_pdf_advanced_python(pptx_path, pdf_path):
        logger.info("✅ Python avanzado exitoso")
        return True
    logger.warning("❌ Python avanzado falló")
    
    logger.error("💥 TODOS LOS MÉTODOS FALLARON")
    return False

# Función auxiliar para instalar dependencias si es necesario
def install_pdf_dependencies():
    """
    Instala las dependencias necesarias para la conversión a PDF
    """
    try:
        import subprocess
        import sys
        
        # Instalar reportlab si no está disponible
        try:
            import reportlab
        except ImportError:
            logger.info("Instalando reportlab...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", "reportlab"])
            
        # Instalar pywin32 si estamos en Windows
        if ON_WINDOWS:
            try:
                import win32com.client
            except ImportError:
                logger.info("Instalando pywin32...")
                subprocess.check_call([sys.executable, "-m", "pip", "install", "pywin32"])
                
    except Exception as e:
        logger.error(f"Error instalando dependencias: {e}")

#  RUTAS

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/procesar', methods=['POST'])
def procesar():
    try:
        if 'excel_file' not in request.files:
            return "No se subió archivo Excel", 400

        excel_file = request.files['excel_file']
        if excel_file.filename == '':
            return "No se seleccionó archivo", 400

        df = pd.read_excel(excel_file)
        logger.info(f"Excel leído correctamente. Filas: {len(df)}")

        # columnas requeridas
        required = ["nombre", "cedula", "horas", "compañia", "certificado"]
        for col in required:
            if col not in [c.lower() for c in df.columns]:
                return f"Columna requerida no encontrada: {col}", 400
        df.columns = [c.lower() for c in df.columns]

        certificados_pendientes = df["certificado"].astype(str).str.lower().str.strip() == "no"
        if not certificados_pendientes.any():
            return render_template("error.html")

        plantilla_path = get_plantilla_path_pptx()
        downloads_folder = get_downloads_folder()
        output_dir = downloads_folder / "Certificados"
        os.makedirs(output_dir, exist_ok=True)

        certificados_por_compania = defaultdict(list)
        certificados_creados = 0

        for index, row in df.iterrows():
            try:
                if str(row["certificado"]).lower().strip() == "no":
                    contexto = {
                        "NOMBRE": str(row["nombre"]),
                        "CEDULA": str(row["cedula"]),
                        "HORAS": str(row["horas"])
                    }

                    compania_folder = output_dir / str(row["compañia"]).replace(" ", "_").replace("/", "_")
                    os.makedirs(compania_folder, exist_ok=True)

                    # Verificar si existe columna 'plantilla', si no usar valor por defecto
                    plantilla_value = row.get('plantilla', 'default') if 'plantilla' in df.columns else 'default'
                    
                    base_name = f"certificado_{plantilla_value}_{row['nombre'].replace(' ', '_')}"
                    base_name = safe_filename(base_name)

                    pptx_file = compania_folder / f"{base_name}.pptx"
                    pdf_file = compania_folder / f"{base_name}.pdf"

                    logger.info(f"Procesando certificado para {row['nombre']}")
                    
                    # Renderizar PPTX
                    render_pptx_template(plantilla_path, contexto, str(pptx_file))
                    if not os.path.exists(pptx_file):
                        logger.error(f"No se pudo crear PPTX: {pptx_file}")
                        continue
                        
                    logger.info(f"PPTX creado exitosamente: {pptx_file}")

                    # Intentar convertir a PDF usando el método definitivo
                    logger.info(f"Iniciando conversión a PDF para: {row['nombre']}")
                    converted = convert_pptx_to_pdf_ultimate(str(pptx_file), str(pdf_file))

                    if converted:
                        # Verificar que el PDF existe y tiene contenido válido
                        if os.path.exists(pdf_file) and os.path.getsize(pdf_file) > 5000:  # Al menos 5KB para un PDF válido
                            try:
                                os.remove(pptx_file)  # Eliminar PPTX solo si PDF es válido
                                logger.info(f"🎉 Certificado PDF creado exitosamente: {pdf_file}")
                                logger.info(f"📄 Tamaño del PDF: {os.path.getsize(pdf_file)} bytes")
                                certificados_por_compania[row["compañia"]].append(os.path.basename(str(pdf_file)))
                            except Exception as e:
                                logger.error(f"Error eliminando PPTX: {e}")
                                certificados_por_compania[row["compañia"]].append(os.path.basename(str(pdf_file)))
                        else:
                            logger.error(f"PDF creado pero parece inválido: {pdf_file}")
                            logger.error(f"Tamaño: {os.path.getsize(pdf_file) if os.path.exists(pdf_file) else 0} bytes")
                            certificados_por_compania[row["compañia"]].append(os.path.basename(str(pptx_file)))
                    else:
                        logger.error(f"❌ No se pudo convertir a PDF: {row['nombre']}")
                        logger.info(f"📋 Manteniendo archivo PPTX: {pptx_file}")
                        certificados_por_compania[row["compañia"]].append(os.path.basename(str(pptx_file)))

                    certificados_creados += 1
                    df.at[index, "certificado"] = "si"

            except Exception as e:
                logger.error(f"Error procesando fila {index} ({row.get('nombre', 'desconocido')}): {e}")
                continue
                

        excel_actualizado = output_dir / "datos_actualizados.xlsx"
        df.to_excel(excel_actualizado, index=False)
        logger.info(f"Excel actualizado guardado: {excel_actualizado}")
        logger.info(f"Procesamiento completado. Certificados creados: {certificados_creados}")

        return render_template("success.html")

    except Exception as e:
        logger.error(f"Error general en procesamiento: {e}")
        return f"Error interno: {str(e)}", 500

@app.errorhandler(500)
def internal_error(error):
    return f"Error interno del servidor: {str(error)}", 500

if __name__ == "__main__":
    threading.Timer(2.0, open_browser).start()
    app.run(host='127.0.0.1', port=5000, debug=False, use_reloader=False)
